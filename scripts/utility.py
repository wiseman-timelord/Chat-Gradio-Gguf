# scripts/utility.py
# v2: Windows 10-11 / Ubuntu 24-25 / Python 3.11-3.13 / Gradio 5.x

# Standard library imports
import re
import subprocess
import json
import time
import psutil
import os
import sys
from pathlib import Path
from datetime import datetime

# Third-party imports
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import gradio as gr
import PyPDF2

# Project imports
import scripts.configure as cfg
from scripts.configure import (
    TEMP_DIR, HISTORY_DIR, SESSION_FILE_FORMAT, ALLOWED_EXTENSIONS,
    current_session_id, session_label
)

# Import search functions from tools module
# NOTE: inference.py is NOT imported here at module level — doing so creates a
# circular dependency (launcher → utility → inference → configure → utility).
# Any inference functions needed inside utility functions must be imported lazily
# inside the function body: `from scripts.inference import func`
from scripts.tools import (
    format_search_status_for_chat,
    web_search, format_web_search_status_for_chat
)

# =============================================================================
# LAZY IMPORTS - spaCy imported only when needed
# =============================================================================
_nlp_model = None

# NOTE: Platform-specific imports (win32com, pythoncom) are imported lazily
# inside the functions that need them, as cfg.PLATFORM is not set until after
# the launcher parses command-line arguments.

def _get_spacy():
    """Lazy import spaCy — deferred to avoid heavy startup cost."""
    try:
        import spacy
        return spacy
    except ImportError:
        return None

# =============================================================================
# BEEP FUNCTION (simple utility, doesn't belong in tools.py)
# =============================================================================

def beep() -> None:
    """Play a notification beep if enabled."""
    if not getattr(cfg, "BLEEP_ON_EVENTS", False):
        return
    if cfg.PLATFORM == "windows":
        _beep_windows()
    elif cfg.PLATFORM == "linux":
        _beep_linux()


def _beep_windows() -> None:
    try:
        import winsound
        winsound.Beep(1000, 150)
    except Exception:
        try:
            import winsound
            winsound.MessageBeep(winsound.MB_OK)
        except Exception:
            pass


def _beep_linux() -> None:
    methods = [
        lambda: subprocess.run(['beep', '-f', '1000', '-l', '150'], timeout=2, check=True,
                               stderr=subprocess.DEVNULL, stdout=subprocess.DEVNULL),
        lambda: subprocess.run(['paplay', '/usr/share/sounds/freedesktop/stereo/complete.oga'],
                               timeout=2, check=True, stderr=subprocess.DEVNULL, stdout=subprocess.DEVNULL)
                if os.path.exists('/usr/share/sounds/freedesktop/stereo/complete.oga')
                else (_ for _ in ()).throw(Exception()),
        lambda: subprocess.run(['play', '-n', 'synth', '0.15', 'sin', '1000'], timeout=2, check=True,
                               stderr=subprocess.DEVNULL, stdout=subprocess.DEVNULL),
        lambda: print("\a", end="", flush=True),
    ]
    for method in methods:
        try:
            method()
            return
        except Exception:
            continue


# =============================================================================
# GENERAL UTILITY FUNCTIONS
# =============================================================================

def short_path(path_str, max_len=44):
    """Truncate path to last max_len chars with ... prefix."""
    path = str(path_str)
    if len(path) <= max_len:
        return path
    return "..." + path[-max_len:]


# =============================================================================
# CPU/GPU DETECTION
# =============================================================================

def detect_cpu_config():
    """Detect CPU configuration and set thread count."""
    try:
        cfg.CPU_PHYSICAL_CORES = psutil.cpu_count(logical=False) or 1
        cfg.CPU_LOGICAL_CORES = psutil.cpu_count(logical=True) or 1

        max_threads = cfg.CPU_LOGICAL_CORES

        if cfg.CPU_THREADS is None or cfg.CPU_THREADS > max_threads:
            cfg.CPU_THREADS = max(1, max_threads // 2)

        if "vulkan" in cfg.BACKEND_TYPE.lower():
            cfg.CPU_THREADS = max(2, cfg.CPU_THREADS)

        print(f"[CPU] Detected: {cfg.CPU_PHYSICAL_CORES} cores, "
              f"{cfg.CPU_LOGICAL_CORES} threads")
        print(f"[CPU] Current: {cfg.CPU_THREADS}")

    except Exception as e:
        cfg.set_status("CPU fallback", console=True)
        print(f"[CPU] Detection error: {e}")
        cfg.CPU_PHYSICAL_CORES = 4
        cfg.CPU_LOGICAL_CORES = 8
        cfg.CPU_THREADS = 4


def get_available_gpus_windows():
    """Retrieve available GPUs on Windows using multiple methods."""
    try:
        output = subprocess.check_output("wmic path win32_VideoController get name", shell=True).decode()
        gpus = [line.strip() for line in output.split('\n') if line.strip() and 'Name' not in line]
        if gpus:
            return gpus
    except Exception:
        pass

    try:
        temp_file = Path(cfg.TEMP_DIR) / "dxdiag.txt"
        subprocess.run(f"dxdiag /t {temp_file}", shell=True, check=True)
        time.sleep(2)
        with open(temp_file, 'r') as f:
            content = f.read()
        gpu = re.search(r"Card name: (.+)", content)
        if gpu:
            return [gpu.group(1).strip()]
    except Exception:
        pass

    return ["CPU Only"]


def get_available_gpus():
    """Get list of available GPUs based on platform."""
    if cfg.PLATFORM == "windows":
        return get_available_gpus_windows()
    else:
        return get_available_gpus_linux()


def get_available_gpus_linux():
    """Retrieve available GPUs on Linux."""
    gpus = []

    try:
        output = subprocess.check_output("lspci | grep -i vga", shell=True).decode()
        for line in output.strip().split('\n'):
            if line:
                parts = line.split(': ')
                if len(parts) > 1:
                    gpus.append(parts[1].strip()[:50])
    except Exception:
        pass

    if not gpus:
        try:
            output = subprocess.check_output(
                "vulkaninfo --summary 2>/dev/null | grep deviceName", shell=True
            ).decode()
            for line in output.strip().split('\n'):
                if 'deviceName' in line:
                    name = line.split('=')[1].strip() if '=' in line else line.split(':')[1].strip()
                    gpus.append(name[:50])
        except Exception:
            pass

    return gpus if gpus else ["CPU Only"]


def get_cpu_info():
    """Get CPU information for display."""
    try:
        cpu_info = []

        if cfg.PLATFORM == "windows":
            try:
                output = subprocess.check_output("wmic cpu get name", shell=True).decode()
                for line in output.split('\n'):
                    line = line.strip()
                    if line and 'Name' not in line:
                        cpu_info.append({
                            "label": line[:50],
                            "cores": cfg.CPU_PHYSICAL_CORES,
                            "threads": cfg.CPU_LOGICAL_CORES
                        })
            except Exception:
                pass
        else:
            try:
                with open('/proc/cpuinfo', 'r') as f:
                    content = f.read()
                model_match = re.search(r'model name\s*:\s*(.+)', content)
                if model_match:
                    cpu_info.append({
                        "label": model_match.group(1).strip()[:50],
                        "cores": cfg.CPU_PHYSICAL_CORES,
                        "threads": cfg.CPU_LOGICAL_CORES
                    })
            except Exception:
                pass

        if not cpu_info:
            cpu_info.append({
                "label": f"CPU ({cfg.CPU_PHYSICAL_CORES}c/{cfg.CPU_LOGICAL_CORES}t)",
                "cores": cfg.CPU_PHYSICAL_CORES,
                "threads": cfg.CPU_LOGICAL_CORES
            })

        return cpu_info

    except Exception as e:
        print(f"[CPU-INFO] Error: {e}")
        return [{"label": "Default CPU", "cores": 4, "threads": 8}]


# =============================================================================
# SESSION MANAGEMENT
# =============================================================================

def get_nlp_model():
    """Get or initialize the spaCy NLP model (lazy load)."""
    global _nlp_model

    spacy = _get_spacy()
    if spacy is None:
        return None  # Graceful fallback if spacy not installed

    if _nlp_model is None:
        try:
            _nlp_model = spacy.load("en_core_web_sm")
        except OSError:
            try:
                subprocess.run([sys.executable, "-m", "spacy", "download", "en_core_web_sm"], check=True)
                _nlp_model = spacy.load("en_core_web_sm")
            except Exception as e:
                print(f"[NLP] Failed to load spaCy model: {e}")
                return None

    return _nlp_model


def summarize_session(messages):
    """Generate a short label for a session based on initial messages."""
    if not messages:
        return "Empty Session"

    first_user_msg = None
    for msg in messages:
        if msg.get('role') == 'user':
            first_user_msg = msg.get('content', '')
            break

    if not first_user_msg:
        return "Session " + datetime.now().strftime("%Y%m%d_%H%M")

    text = first_user_msg[:500]

    nlp = get_nlp_model()
    if nlp:
        try:
            doc = nlp(text)
            candidates = []
            for ent in doc.ents:
                if ent.label_ in ['PERSON', 'ORG', 'PRODUCT', 'EVENT', 'WORK_OF_ART']:
                    candidates.append(ent.text)
            for chunk in doc.noun_chunks:
                if chunk.root.pos_ in ['NOUN', 'PROPN']:
                    candidates.append(chunk.text)

            if candidates:
                label = candidates[0][:40]
                return label.strip()
        except Exception:
            pass

    words = text.split()[:6]
    return ' '.join(words)[:40] if words else "Session"


def get_saved_sessions():
    """Get list of saved session files, sorted by date (newest first)."""
    history_path = Path(HISTORY_DIR)
    if not history_path.exists():
        return []

    sessions = []
    for file in history_path.glob("*.json"):
        try:
            mtime = file.stat().st_mtime
            sessions.append((mtime, file.name))
        except Exception:
            continue

    sessions.sort(reverse=True)
    return [s[1] for s in sessions[:cfg.MAX_HISTORY_SLOTS]]


def save_session_history(messages, attached_files=None):
    """Save current session to history."""
    if not messages:
        return

    session_id = cfg.current_session_id or datetime.now().strftime("%Y%m%d_%H%M%S")
    label = cfg.session_label or summarize_session(messages)

    cfg.current_session_id = session_id
    cfg.session_label = label

    filename = f"{session_id}.json"
    filepath = Path(HISTORY_DIR) / filename

    session_data = {
        "session_id": session_id,
        "label": label,
        "timestamp": datetime.now().isoformat(),
        "messages": messages,
        "attached_files": attached_files or []
    }

    try:
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(session_data, f, indent=2, ensure_ascii=False)
        print(f"[SESSION] Saved: {label} ({filename})")
    except Exception as e:
        print(f"[SESSION] Save error: {e}")


def load_session_history(filename):
    """Load a session from history file."""
    filepath = Path(HISTORY_DIR) / filename

    if not filepath.exists():
        return None, None, [], []

    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)

        session_id = data.get("session_id", "")
        label = data.get("label", "Unnamed Session")
        messages = data.get("messages", [])
        attached_files = data.get("attached_files", [])

        return session_id, label, messages, attached_files

    except Exception as e:
        print(f"[SESSION] Load error: {e}")
        return None, None, [], []


# =============================================================================
# FILE HANDLING
# =============================================================================

def read_file_content(file_path):
    """Read content from various file types."""
    try:
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
            encodings_to_try = ['utf-8', 'cp1252', 'iso-8859-1', 'latin-1']
            content = None
            for encoding in encodings_to_try:
                try:
                    with open(path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if content is not None:
                return content, "text", True, None
            else:
                with open(path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read(), "text", True, None

        elif suffix == '.pdf':
            reader = PyPDF2.PdfReader(str(path))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return text, "text", True, None

        elif suffix == '.docx':
            doc = Document(str(path))
            text = "\n".join(para.text for para in doc.paragraphs)
            return text, "text", True, None

        elif suffix == '.xlsx':
            wb = load_workbook(str(path), data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            return "\n".join(text_parts), "text", True, None

        elif suffix == '.pptx':
            prs = Presentation(str(path))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts), "text", True, None

        elif suffix in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
            import base64
            with open(path, 'rb') as f:
                data = base64.b64encode(f.read()).decode('utf-8')
            mime_type = {
                '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp'
            }.get(suffix, 'image/png')
            return f"data:{mime_type};base64,{data}", "image", True, None

        else:
            return None, None, False, f"Unsupported file type: {suffix}"

    except Exception as e:
        return None, None, False, str(e)


def eject_file(file_list, slot_index, is_attach=True):
    """Eject a file from the specified slot."""
    if 0 <= slot_index < len(file_list):
        removed_file = file_list.pop(slot_index)
        if is_attach:
            cfg.session_attached_files = file_list
        status_msg = f"Ejected {Path(removed_file).name}"
        print(f"[FILES] {status_msg}")
    else:
        status_msg = "No file to eject"

    return file_list, status_msg


def update_file_slot_ui(file_list, is_attach=True):
    """Update file slot UI components."""
    max_slots = cfg.MAX_POSSIBLE_ATTACH_SLOTS
    current_max = cfg.MAX_ATTACH_SLOTS
    button_updates = []

    for i in range(max_slots):
        if i >= current_max:
            button_updates.append(gr.update(value="", visible=False, variant="primary"))
        elif i < len(file_list):
            filename = Path(file_list[i]).name
            short_name = filename[:36] + ".." if len(filename) > 38 else filename
            button_updates.append(gr.update(value=short_name, visible=True, variant="primary"))
        else:
            button_updates.append(gr.update(value="", visible=False, variant="primary"))

    show_upload = len(file_list) < current_max if is_attach else True
    button_updates.append(gr.update(visible=show_upload))

    return button_updates


def process_attach_files(files, attached_files):
    """Process uploaded files for attachment."""
    if not files:
        return "No files selected.", attached_files
    status, updated_files = process_files(files, attached_files, cfg.MAX_ATTACH_SLOTS, is_attach=True)
    cfg.session_attached_files = updated_files
    return status, updated_files


def process_files(files, existing_files, max_files, is_attach=True):
    """Process uploaded files, ensuring no duplicates and respecting max file limits."""
    print(f"[FILES] Raw input type: {type(files)}, value: {files}")

    if not files:
        return "No files uploaded.", existing_files

    # Handle Gradio 5.x file upload format — files can be list of file objects or dicts
    normalized_files = []
    if isinstance(files, (list, tuple)):
        for f in files:
            if f is None:
                continue
            if hasattr(f, 'name'):
                file_path = f.name
            elif isinstance(f, (str, Path)):
                file_path = str(f)
            elif isinstance(f, dict):
                file_path = f.get('name') or f.get('path') or f.get('data')
            else:
                print(f"[FILES] Skipping unknown file format: {type(f)} - {f}")
                continue

            if file_path and os.path.isfile(file_path):
                normalized_files.append(str(file_path))
                print(f"[FILES] Added file: {file_path}")
            else:
                print(f"[FILES] File not found or invalid: {file_path}")
    elif isinstance(files, str):
        if os.path.isfile(files):
            normalized_files.append(files)
    elif hasattr(files, 'name'):
        normalized_files.append(files.name)

    print(f"[FILES] Normalized {len(normalized_files)} files")

    if not normalized_files:
        return "No valid files to add.", existing_files

    new_files = [f for f in normalized_files if f not in existing_files]
    if not new_files:
        return "No new files to add.", existing_files

    # Remove any existing files with the same filename as the incoming ones
    for f in new_files:
        file_name = Path(f).name
        existing_files = [ef for ef in existing_files if Path(ef).name != file_name]

    available_slots = max_files - len(existing_files)
    processed_files = new_files[:available_slots]
    updated_files = processed_files + existing_files

    if is_attach:
        cfg.session_attached_files = updated_files
    else:
        cfg.session_vector_files = updated_files

    status = f"Attached {len(processed_files)} file(s)."
    print(f"[FILES] Status: {status}, Total files: {len(updated_files)}")
    return status, updated_files