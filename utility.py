# Script: `.\scripts\utility.py`

# Imports
import tempfile
import re, subprocess, json, time, random, psutil, shutil, os, zipfile, spacy, sys, PyPDF2
from pathlib import Path
from datetime import datetime
from newspaper import Article
from ddgs import DDGS
from ddgs.exceptions import DDGSException
import requests.exceptions  # For HTTPError and Timeout
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from .models import load_models, clean_content
import scripts.settings as settings
from .temporary import (
    TEMP_DIR, HISTORY_DIR, SESSION_FILE_FORMAT, ALLOWED_EXTENSIONS, 
    current_session_id, session_label
)
from . import temporary


# Variables
_nlp_model = None

# Conditional imports based on platform
if temporary.PLATFORM == "windows":
    import win32com.client
    import pythoncom
elif temporary.PLATFORM == "linux":
    try:
        import pyttsx3
    except ImportError:
        print("Warning: pyttsx3 not installed. Text-to-speech will be unavailable on Linux.")

# Functions...
def short_path(path_str, max_len=44):
    """Truncate path to last max_len chars with ... prefix"""
    path = str(path_str)
    if len(path) <= max_len:
        return path
    return "..." + path[-max_len:]

def filter_operational_content(text):
    """Remove operational tags and metadata from the text."""
    text = re.sub(r'^AI-Chat:\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    text = re.sub(r'<answer>.*?</answer>', '', text, flags=re.DOTALL)
    return text.strip()

def detect_cpu_config():
    """Detect CPU configuration and set thread options"""
    try:
        import psutil
        cpu_info = {
            "physical_cores": psutil.cpu_count(logical=False) or 1,
            "logical_cores": psutil.cpu_count(logical=True) or 1
        }
        
        temporary.CPU_PHYSICAL_CORES = cpu_info["physical_cores"]
        temporary.CPU_LOGICAL_CORES = cpu_info["logical_cores"]
        
        # Generate thread options (1 to logical_cores, not logical_cores-1)
        # This allows full utilization of available threads
        max_threads = temporary.CPU_LOGICAL_CORES
        temporary.CPU_THREAD_OPTIONS = list(range(1, max_threads + 1))
        
        # Set default threads if not already configured
        if temporary.CPU_THREADS is None or temporary.CPU_THREADS > max_threads:
            # Default to 50% of available threads, minimum 1
            temporary.CPU_THREADS = max(1, max_threads // 2)
        
        # Vulkan-specific: Even with all layers on GPU, still uses some CPU threads
        if "vulkan" in temporary.BACKEND_TYPE.lower():
            temporary.CPU_THREADS = max(2, temporary.CPU_THREADS)
            
        print(f"[CPU] Detected: {temporary.CPU_PHYSICAL_CORES} cores, "
              f"{temporary.CPU_LOGICAL_CORES} threads")
        print(f"[CPU] Available thread range: 1-{max_threads}, "
              f"Current: {temporary.CPU_THREADS}")
            
    except Exception as e:
        temporary.set_status("CPU fallback", console=True)
        print(f"[CPU] Detection error: {e}")
        # Fallback values
        temporary.CPU_PHYSICAL_CORES = 4
        temporary.CPU_LOGICAL_CORES = 8
        temporary.CPU_THREAD_OPTIONS = [1, 2, 3, 4, 5, 6, 7, 8]
        temporary.CPU_THREADS = 4

def get_available_gpus_windows():
    """
    Retrieve available GPUs on Windows using wmic and dxdiag.
    """
    try:
        output = subprocess.check_output("wmic path win32_VideoController get name", shell=True).decode()
        gpus = [line.strip() for line in output.split('\n') if line.strip() and 'Name' not in line]
        return gpus if gpus else ["CPU Only"]
    except Exception:
        try:
            temp_file = Path(temporary.TEMP_DIR) / "dxdiag.txt"
            subprocess.run(f"dxdiag /t {temp_file}", shell=True, check=True)
            time.sleep(2)
            with open(temp_file, 'r') as f:
                content = f.read()
                gpu = re.search(r"Card name: (.+)", content)
                return [gpu.group(1).strip()] if gpu else ["CPU Only"]
        except Exception:
            return ["CPU Only"]

def get_available_gpus_linux():
    """Get available GPUs on Linux systems with improved detection"""
    gpus = []

    try:
        # 1) NVIDIA GPUs (via nvidia-smi)
        try:
            output = subprocess.check_output(
                ["nvidia-smi", "--query-gpu=name", "--format=csv,noheader"],
                stderr=subprocess.DEVNULL,
                text=True
            )
            gpus.extend(f"NVIDIA {line.strip()}" for line in output.splitlines() if line.strip())
        except Exception:
            pass

        # 2) AMD GPUs (via rocminfo – ROCm stack)
        try:
            output = subprocess.check_output(
                ["rocminfo"],
                stderr=subprocess.DEVNULL,
                text=True
            )
            for line in output.splitlines():
                if "Marketing Name" in line:
                    name = line.split(":", 1)[-1].strip()
                    if name:
                        gpus.append(f"AMD {name}")
        except Exception:
            pass

        # 3) Generic PCI bus scan (catches AMD, Intel, and others)
        try:
            output = subprocess.check_output(
                ["lspci", "-nn"],
                stderr=subprocess.DEVNULL,
                text=True
            )
            for line in output.splitlines():
                lower = line.lower()
                if "vga" in lower or "display" in lower or "3d" in lower:
                    name = line.split(":", 2)[-1].strip()
                    if "nvidia" in lower:
                        gpus.append(f"NVIDIA {name}")
                    elif "amd" in lower or "ati" in lower:
                        gpus.append(f"AMD {name}")
                    elif "intel" in lower:
                        gpus.append(f"Intel {name}")
                    else:
                        gpus.append(name)
        except Exception:
            pass

    except Exception as e:
        print(f"GPU detection error: {str(e)}")

    # Remove duplicates while preserving order
    seen = set()
    unique_gpus = [g for g in gpus if not (g in seen or seen.add(g))]

    if not unique_gpus:
        temporary.set_status("No GPU", console=True)
        sys.exit(1)          # critical exit → back to launcher caller

    return unique_gpus

def get_cpu_info():
    """
    Get CPU information for configuration purposes.
    Returns a list of dictionaries with CPU information.
    """
    try:
        import psutil
        cpu_count = psutil.cpu_count(logical=False) or 1
        logical_count = psutil.cpu_count(logical=True) or 1
        
        # Try to get CPU brand/model info
        try:
            with open('/proc/cpuinfo', 'r') as f:
                cpuinfo = f.read()
            model_match = re.search(r'model name\s*:\s*(.+)', cpuinfo)
            if model_match:
                model = model_match.group(1).strip()
            else:
                model = "Unknown CPU"
        except:
            model = "Generic CPU"
        
        return [{
            "label": f"{model} ({cpu_count} cores, {logical_count} threads)",
            "physical_cores": cpu_count,
            "logical_cores": logical_count
        }]
    except ImportError:
        # Fallback if psutil is not available
        return [{
            "label": "Generic CPU",
            "physical_cores": 4,
            "logical_cores": 8
        }]
    except Exception as e:
        print(f"Error getting CPU info: {e}")
        return [{
            "label": "Default CPU",
            "physical_cores": 4,
            "logical_cores": 8
        }]

def get_available_gpus():
    """Returns list of GPUs with Intel marked appropriately"""
    gpus = []
    if temporary.PLATFORM == "windows":
        # Windows detection logic that identifies Intel GPUs
        gpus = get_available_gpus_windows()  # Should include Intel GPUs
    elif temporary.PLATFORM == "linux":
        # Linux detection that identifies Intel GPUs
        gpus = get_available_gpus_linux()    # Should include Intel GPUs
    return gpus
            
def generate_session_id():
    return datetime.now().strftime("%Y%m%d_%H%M%S")

def speak_text(text):
    """Speak the given text using platform-specific text-to-speech with improved Windows support."""
    if not text or not isinstance(text, str):
        return
    
    # Clean the text
    text = text.strip()
    if not text:
        return
        
    # Check if we just spoke this exact text (optional - can be removed if too restrictive)
    if hasattr(temporary, 'LAST_SPOKEN') and temporary.LAST_SPOKEN == text:
        print("[TTS] Skipping duplicate speech")
        return
        
    if temporary.PLATFORM == "windows":
        import threading
        
        def _speak_windows_threaded():
            """Windows TTS in dedicated thread with proper COM initialization."""
            import pythoncom
            import win32com.client
            
            # Initialize COM with Single-Threaded Apartment (STA) - CRITICAL for Windows 8.1
            try:
                pythoncom.CoInitialize()  # STA initialization
                
                try:
                    # Create SAPI voice object
                    speaker = win32com.client.Dispatch("SAPI.SpVoice")
                    
                    # Optional: Configure voice properties
                    # speaker.Rate = 0  # -10 to 10, 0 is default
                    # speaker.Volume = 100  # 0 to 100
                    
                    # Speak synchronously (blocks until complete)
                    speaker.Speak(text)
                    
                    # Mark as spoken
                    temporary.LAST_SPOKEN = text
                    print(f"[TTS-WIN] Spoke: {text[:50]}...")
                    
                except Exception as e:
                    print(f"[TTS-WIN] SAPI error: {str(e)}")
                    # Fallback beep to indicate speech failure
                    try:
                        import winsound
                        winsound.Beep(1000, 200)
                    except:
                        pass
                        
            finally:
                # CRITICAL: Always uninitialize COM
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
        
        # Run speech in dedicated thread to avoid COM/async conflicts
        thread = threading.Thread(target=_speak_windows_threaded, daemon=True)
        thread.start()
        thread.join(timeout=30)  # Wait up to 30 seconds
        
        if thread.is_alive():
            print("[TTS-WIN] Speech timed out")
                
    elif temporary.PLATFORM == "linux":
        try:
            # Initialize engine if not exists
            if not hasattr(temporary, 'tts_engine'):
                import pyttsx3
                try:
                    temporary.tts_engine = pyttsx3.init(driverName='espeak')
                    temporary.tts_engine.setProperty('rate', 150)
                    temporary.tts_engine.setProperty('volume', 0.9)
                except:
                    temporary.tts_engine = pyttsx3.init()
                    
            temporary.tts_engine.say(text)
            temporary.tts_engine.runAndWait()
            temporary.LAST_SPOKEN = text
            print(f"[TTS-LINUX] Spoke: {text[:50]}...")
            
        except Exception as e:
            print(f"[TTS-LINUX] pyttsx3 error: {str(e)}")
            try:
                # Fallback to direct espeak
                subprocess.run(
                    ['espeak', '-v', 'en-us', '-s', '150', text],
                    check=False,
                    timeout=30
                )
                print(f"[TTS-LINUX] Fallback espeak: {text[:50]}...")
            except FileNotFoundError:
                try:
                    # Final fallback to spd-say
                    subprocess.run(
                        ['spd-say', '--wait', '-r', '-50', '-t', 'female3', text],
                        check=False,
                        timeout=30
                    )
                    print(f"[TTS-LINUX] Fallback spd-say: {text[:50]}...")
                except FileNotFoundError:
                    print("[TTS-LINUX] All TTS methods unavailable")
                except Exception as e:
                    print(f"[TTS-LINUX] spd-say error: {str(e)}")
            except Exception as e:
                print(f"[TTS-LINUX] espeak error: {str(e)}")
    else:
        raise ValueError(f"Unsupported platform: {temporary.PLATFORM}")

def get_nlp_model():
    """Lazy load spaCy model on first use."""
    global _nlp_model
    if _nlp_model is None:
        try:
            _nlp_model = spacy.load("en_core_web_sm")
            print("[SPACY] Language model loaded")
        except OSError:
            print("[SPACY] WARNING: Model not found. Run: python -m spacy download en_core_web_sm")
            _nlp_model = False  # Mark as failed to avoid repeated attempts
    return _nlp_model if _nlp_model is not False else None

def read_file_content(file_path, max_chars=50000):
    """
    Read file content with support for multiple formats.
    Returns tuple: (content: str, file_type: str, success: bool, error: str)
    
    Supported formats:
    - Text: .txt, .py, .json, .yaml, .md, .xml, .html, .css, .js, .sh, .bat, .ps1
    - PDF: .pdf
    - Word: .docx
    - Excel: .xlsx, .xls
    - PowerPoint: .pptx
    """
    from pathlib import Path
    
    file_path = Path(file_path)
    if not file_path.exists():
        return "", "unknown", False, f"File not found: {file_path}"
    
    extension = file_path.suffix.lower()
    
    try:
        # Text-based files
        text_extensions = {'.txt', '.py', '.json', '.yaml', '.yml', '.md', 
                          '.xml', '.html', '.css', '.js', '.sh', '.bat', 
                          '.ps1', '.psd1', '.xaml', '.csv', '.log'}
        
        if extension in text_extensions:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(max_chars)
                return content, "text", True, ""
        
        # PDF files
        elif extension == '.pdf':
            try:
                import PyPDF2
                content = []
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    max_pages = min(10, len(reader.pages))  # Limit to 10 pages
                    for page_num in range(max_pages):
                        page = reader.pages[page_num]
                        content.append(page.extract_text())
                text = '\n'.join(content)[:max_chars]
                return text, "pdf", True, ""
            except ImportError:
                return "", "pdf", False, "PyPDF2 not installed. Run: pip install PyPDF2"
            except Exception as e:
                return "", "pdf", False, f"PDF error: {str(e)}"
        
        # Word documents
        elif extension == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                content = '\n'.join([para.text for para in doc.paragraphs])[:max_chars]
                return content, "docx", True, ""
            except ImportError:
                return "", "docx", False, "python-docx not installed. Run: pip install python-docx"
            except Exception as e:
                return "", "docx", False, f"DOCX error: {str(e)}"
        
        # Excel files
        elif extension in {'.xlsx', '.xls'}:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, read_only=True, data_only=True)
                content = []
                for sheet_name in wb.sheetnames[:3]:  # Limit to 3 sheets
                    sheet = wb[sheet_name]
                    content.append(f"=== Sheet: {sheet_name} ===")
                    for row in list(sheet.rows)[:50]:  # Limit to 50 rows per sheet
                        row_data = [str(cell.value) if cell.value is not None else "" 
                                   for cell in row]
                        content.append(" | ".join(row_data))
                text = '\n'.join(content)[:max_chars]
                return text, "excel", True, ""
            except ImportError:
                return "", "excel", False, "openpyxl not installed. Run: pip install openpyxl"
            except Exception as e:
                return "", "excel", False, f"Excel error: {str(e)}"
        
        # PowerPoint files
        elif extension == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                content = []
                for i, slide in enumerate(prs.slides[:10], 1):  # Limit to 10 slides
                    content.append(f"=== Slide {i} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content.append(shape.text)
                text = '\n'.join(content)[:max_chars]
                return text, "pptx", True, ""
            except ImportError:
                return "", "pptx", False, "python-pptx not installed. Run: pip install python-pptx"
            except Exception as e:
                return "", "pptx", False, f"PowerPoint error: {str(e)}"
        
        else:
            return "", "unsupported", False, f"Unsupported file type: {extension}"
    
    except Exception as e:
        return "", "error", False, f"Unexpected error: {str(e)}"

def get_attached_files_context(attached_files, query=None, max_total_chars=8000):
    """
    Generate context string from attached files for model input.
    
    Args:
        attached_files: List of file paths
        query: Optional user query for RAG-based filtering
        max_total_chars: Maximum total characters to include
    
    Returns:
        str: Formatted context string with file contents
    """
    if not attached_files:
        return ""
    
    from pathlib import Path
    import scripts.temporary as temporary
    
    context_parts = []
    total_chars = 0
    
    # Build list of files with their names
    file_list = [Path(f).name for f in attached_files]
    context_parts.append(f"📎 Attached Files ({len(file_list)}): {', '.join(file_list)}")
    context_parts.append("")
    
    # Try RAG-based retrieval first if query provided
    rag_context = None
    if query and hasattr(temporary.context_injector, 'get_relevant_context'):
        try:
            rag_context = temporary.context_injector.get_relevant_context(query, k=3)
            if rag_context:
                context_parts.append("📖 Relevant Content from Files:")
                context_parts.append(rag_context[:max_total_chars // 2])
                context_parts.append("")
                total_chars += len(rag_context)
        except Exception as e:
            print(f"[ATTACH] RAG retrieval error: {e}")
    
    # For small files, include full content
    remaining_chars = max_total_chars - total_chars
    if remaining_chars > 1000:
        for file_path in attached_files:
            if total_chars >= max_total_chars:
                context_parts.append(f"... ({len(attached_files) - len(context_parts) + 3} more files)")
                break
            
            content, file_type, success, error = read_file_content(
                file_path, 
                max_chars=remaining_chars // len(attached_files)
            )
            
            file_name = Path(file_path).name
            
            if success and content.strip():
                # Skip if content already in RAG context
                if rag_context and content[:500] in rag_context:
                    continue
                    
                context_parts.append(f"--- {file_name} ({file_type}) ---")
                context_parts.append(content)
                context_parts.append("")
                total_chars += len(content)
                remaining_chars -= len(content)
            elif not success:
                context_parts.append(f"⚠️ {file_name}: {error}")
    
    return "\n".join(context_parts)

def sanitize_label(label: str) -> str:
    """
    Make a session-label JSON-safe and Windows-cp1252-safe.
    Returns 'Untitled' if nothing usable remains.
    """
    if not label:
        return "Untitled"

    # 1. Remove C0 / C1 control codes (0x00-0x1F, 0x7F-0x9F) incl. 0x9d
    label = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', label)

    # 2. Replace common Unicode punctuation with ASCII
    label = (label
             .replace('\u2018', "'").replace('\u2019', "'")   # ‘ ’
             .replace('\u201C', '"').replace('\u201D', '"')   # “ ”
             .replace('\u2013', '-').replace('\u2014', '-')   # – —
             .replace('\u2026', '...')                        # …
             .replace('\u00A0', ' '))                         # non-breaking space

    # 3. Drop any remaining non-printable / non-basic-plane code-points
    label = ''.join(c for c in label if c.isprintable() and ord(c) < 0x10000)

    # 4. Collapse multiple spaces / leading-trailing junk
    label = re.sub(r'\s+', ' ', label).strip()

    # 5. Final guard
    if not label:
        return "Untitled"

    # 6. Length cap
    return label[:50].strip()

# Replace generate_session_label function:
def generate_session_label(session_log):
    """
    Generate a session label using spaCy keyword extraction, up to 50 characters.
    Output is sanitised for JSON safety.
    """
    if not session_log:
        return "Untitled"

    nlp = get_nlp_model()
    if nlp is None:
        # fallback: first few words of first user message
        for msg in session_log:
            if msg['role'] == 'user':
                content = clean_content(msg['role'], msg['content'])
                words = content.split()[:3]
                label = ' '.join(words) if words else "Untitled"
                return sanitize_label(label)
        return "Untitled"

    text_for_analysis = " ".join(
        clean_content(msg['role'], msg['content']) for msg in session_log
    )
    text_for_analysis = filter_operational_content(text_for_analysis)
    if not text_for_analysis.strip():
        return "Untitled"

    try:
        doc = nlp(text_for_analysis[:1000])   # speed guard
        candidates = []

        for ent in doc.ents:
            if ent.label_ in {'PERSON', 'ORG', 'GPE', 'PRODUCT', 'EVENT'}:
                candidates.append((ent.text, len(ent.text.split())))

        for chunk in doc.noun_chunks:
            if chunk.root.pos_ in {'NOUN', 'PROPN'}:
                candidates.append((chunk.text, len(chunk.text.split())))

        for token in doc:
            if token.pos_ in {'NOUN', 'PROPN'} and not token.is_stop:
                candidates.append((token.text, 1))

        if candidates:
            candidates.sort(key=lambda x: (-x[1], text_for_analysis.find(x[0])))
            description = candidates[0][0]
        else:
            words = [t.text for t in doc if not t.is_stop and t.is_alpha][:3]
            description = ' '.join(words) if words else "No description"

        return sanitize_label(description)          # ← critical
    except Exception as e:
        print(f"[SPACY] label error: {e}")
        words = text_for_analysis.split()[:3]
        label = ' '.join(words) if words else "Untitled"
        return sanitize_label(label)

def save_session_history(session_log, attached_files):
    """
    Save session history with UTF-8 encoding.
    Preserves the current label if already present (prevents re-editing corruption).
    """
    try:
        Path(HISTORY_DIR).mkdir(parents=True, exist_ok=True)
        Path(TEMP_DIR).mkdir(parents=True, exist_ok=True)

        if not temporary.current_session_id:
            temporary.current_session_id = generate_session_id()

        # ← do NOT regenerate label if one exists
        if not temporary.session_label or temporary.session_label == "Untitled":
            temporary.session_label = generate_session_label(session_log)

        session_file = Path(HISTORY_DIR) / f"session_{temporary.current_session_id}.json"
        session_data = {
            "session_id": temporary.current_session_id,
            "label": temporary.session_label,
            "history": session_log,
            "attached_files": [str(Path(f).resolve()) for f in attached_files] if attached_files else [],
            "last_saved": datetime.now().isoformat()
        }

        temp_file = Path(TEMP_DIR) / f"temp_{temporary.current_session_id}.json"
        with open(temp_file, 'w', encoding='utf-8') as f:   # ← UTF-8 enforced
            json.dump(session_data, f, ensure_ascii=False, indent=2)

        temp_file.replace(session_file)
        temporary.set_status("Ready")
        manage_session_history()
        return True
    except Exception as e:
        print(f"Error saving session: {e}")
        import traceback
        traceback.print_exc()
        return False

def load_session_history(session_file):
    try:
        with open(session_file, 'r', encoding='utf-8', errors='replace') as f:
            data = json.load(f)
    except UnicodeDecodeError as e:
        print(f"Encoding error in {session_file}: {e}")
        try:
            with open(session_file, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)
            print("Recovered session with ignored errors")
        except Exception as e2:
            print(f"Failed to recover {session_file}: {e2}")
            return None, "Corrupted", [], []
    except Exception as e:
        print(f"Error loading session file {session_file}: {e}")
        return None, "Error", [], []

    session_id = data.get("session_id", session_file.stem.replace('session_', ''))
    label = sanitize_label(data.get("label", "Untitled"))
    history = data.get("history", [])
    attached_files = [str(Path(f).resolve()) for f in data.get("attached_files", []) if Path(f).exists()]

    if len(attached_files) != len(data.get("attached_files", [])):
        print(f"Removed missing attached files from session {session_id}")

    temporary.session_attached_files = attached_files
    return session_id, label, history, attached_files

def manage_session_history():
    """Limit saved sessions to MAX_HISTORY_SLOTS."""
    history_dir = Path(HISTORY_DIR)
    session_files = sorted(history_dir.glob("session_*.json"), key=lambda x: x.stat().st_mtime, reverse=True)
    while len(session_files) > temporary.MAX_HISTORY_SLOTS:
        oldest_file = session_files.pop()
        oldest_file.unlink()
        print(f"Deleted oldest session: {oldest_file}")

def process_uploaded_files(files):
    """Process uploaded files (if needed beyond copying)."""
    return [file.name for file in files] if files else []

def chunk_text_for_speech(text, max_chars=500):
    chunks = []
    current_chunk = []
    current_length = 0
    
    # Split into sentences first
    sentences = re.split(r'(?<=[.!?])\s+', text)
    
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        sentence_length = len(sentence)
        
        if current_length + sentence_length <= max_chars:
            current_chunk.append(sentence)
            current_length += sentence_length
        else:
            if current_chunk:
                chunks.append(' '.join(current_chunk))
                current_chunk = []
                current_length = 0
            # Handle long individual sentences
            while sentence_length > max_chars:
                split_pos = sentence[:max_chars].rfind('. ') + 1
                if split_pos <= 0:
                    split_pos = max_chars
                chunks.append(sentence[:split_pos])
                sentence = sentence[split_pos:].lstrip()
                sentence_length = len(sentence)
            if sentence:
                current_chunk.append(sentence)
                current_length += sentence_length
                
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    return chunks

def web_search(query: str, num_results, max_hits: int = 6) -> str:
    """
    DuckDuckGo text search that also dumps raw results to terminal
    when temporary.PRINT_RAW_OUTPUT is True.
    """
    import scripts.temporary as tmp           # local import avoids cycles

    if not query.strip():
        return "Empty query."

    try:
        hits = DDGS().text(query, max_results=max_hits)
    except DDGSException as e:
        raw = f"DuckDuckGo error: {e}"
        if tmp.PRINT_RAW_OUTPUT:
            print("=== RAW DDG ===\n", raw, "\n=== END ===", flush=True)
        return raw

    if not hits:
        raw = "DuckDuckGo returned zero results."
        if tmp.PRINT_RAW_OUTPUT:
            print("=== RAW DDG ===\n", raw, "\n=== END ===", flush=True)
        return raw

    raw = "\n\n".join(
        f"[{i}] **{h.get('title','')}**\n{h.get('body','')}\n*Source:* <{h.get('href','')}>"
        for i, h in enumerate(hits, 1)
    )

    if tmp.PRINT_RAW_OUTPUT:
        print("=== RAW DDG ===\n", raw, "\n=== END ===", flush=True)

    return raw

def summarize_document(file_path):
    """Summarize the contents of a document using spaCy, up to 100 characters."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        nlp = get_nlp_model()
        if nlp is None:
            # Fallback: first few words
            words = content.split()[:10]
            return ' '.join(words)[:100] if words else "No summary available"
        
        # Process with spaCy
        doc = nlp(content[:2000])  # Limit to first 2000 chars
        
        # Extract key phrases
        candidates = []
        
        # Get entities
        for ent in doc.ents:
            candidates.append(ent.text)
        
        # Get noun chunks
        for chunk in doc.noun_chunks:
            if chunk.root.pos_ in ['NOUN', 'PROPN']:
                candidates.append(chunk.text)
        
        if candidates:
            summary = candidates[0]
        else:
            # Fallback to first sentence
            sentences = [sent.text.strip() for sent in doc.sents]
            summary = sentences[0] if sentences else "No summary available"
        
        return summary[:100]
        
    except Exception as e:
        print(f"Error summarizing document {file_path}: {e}")
        return "Error generating summary"

def get_attached_files_summary(attached_files):
    """Generate a list of summaries for attached files."""
    if not attached_files:
        return "No attached files to summarize."
    summary_list = []
    for file in attached_files:
        summary = summarize_document(file)
        summary_list.append(f"**{Path(file).name}** - {summary}")
    return "\n".join(summary_list)

def load_and_chunk_documents(file_paths: list) -> list:
    """Load and chunk documents from a list of file paths for RAG."""
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain_community.document_loaders import TextLoader
    from .temporary import CONTEXT_SIZE, RAG_CHUNK_SIZE_DEVIDER, RAG_CHUNK_OVERLAP_DEVIDER
    documents = []
    try:
        chunk_size = CONTEXT_SIZE // (RAG_CHUNK_SIZE_DEVIDER if RAG_CHUNK_SIZE_DEVIDER != 0 else 4)
        chunk_overlap = CONTEXT_SIZE // (RAG_CHUNK_OVERLAP_DEVIDER if RAG_CHUNK_OVERLAP_DEVIDER != 0 else 32)
        for file_path in file_paths:
            if Path(file_path).suffix[1:].lower() in ALLOWED_EXTENSIONS:
                loader = TextLoader(file_path)
                docs = loader.load()
                splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
                chunks = splitter.split_documents(docs)
                documents.extend(chunks)
    except Exception as e:
        print(f"Error loading documents: {e}")
    return documents

def delete_all_session_histories():
    """
    Delete all history JSON files in HISTORY_DIR.
    Returns a status message indicating the result.
    """
    history_dir = Path(HISTORY_DIR)
    for file in history_dir.glob('*.json'):
        try:
            file.unlink()
            print(f"Deleted history file: {file}")
        except Exception as e:
            print(f"Error deleting {file}: {e}")
    return "All session histories deleted."

def get_saved_sessions():
    """Get list of saved session files sorted by modification time."""
    history_dir = Path(HISTORY_DIR)
    session_files = sorted(history_dir.glob("session_*.json"), key=lambda x: x.stat().st_mtime, reverse=True)
    return [f.name for f in session_files]

def create_session_vectorstore(file_paths):
    """
    Thin wrapper so interface.py can call the injector transparently.
    Returns nothing; side-effect is that context_injector now has data.
    """
    from scripts.temporary import context_injector
    context_injector.set_session_vectorstore(file_paths)

def process_files(files, existing_files, max_files, is_attach=True):
    """
    Process uploaded files for attach or vector, ensuring no duplicates and respecting max file limits.
    """
    if not files:
        return "No files uploaded.", existing_files

    new_files = [f for f in files if os.path.isfile(f) and f not in existing_files]
    if not new_files:
        return "No new files to add.", existing_files

    for f in new_files:
        file_name = Path(f).name
        existing_files = [ef for ef in existing_files if Path(ef).name != file_name]

    available_slots = max_files - len(existing_files)
    processed_files = new_files[:available_slots]
    updated_files = processed_files + existing_files

    if is_attach:
        temporary.session_attached_files = updated_files
    else:
        temporary.session_vector_files = updated_files

    status = f"Processed {len(processed_files)} new {'attach' if is_attach else 'vector'} files."
    return status, updated_files

